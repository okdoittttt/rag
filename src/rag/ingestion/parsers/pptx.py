"""PowerPoint(.pptx) 파서

python-pptx를 사용하여 슬라이드의 텍스트와 표를 추출합니다.
"""

from __future__ import annotations

from pathlib import Path

from rag.ingestion.parsers.base import DocumentParser
from rag.logger import get_logger


logger = get_logger(__name__)


class PptxParser(DocumentParser):
    """PowerPoint(.pptx) 파일 파서"""

    extensions = [".pptx"]

    def parse(self, path: Path) -> str:
        """.pptx 파일에서 텍스트를 추출한다.

        슬라이드 순서대로 ``# 슬라이드 N`` 머리말을 붙이고, 각 도형의 텍스트
        프레임과 표 셀 텍스트를 모은다. ``python-pptx`` 는 호출 시점에 지연
        import 하여 선택적 의존성으로 다룬다.

        Args:
            path: 파싱할 .pptx 파일 경로.

        Returns:
            추출된 텍스트. 추출 실패 시 빈 문자열.

        Raises:
            ValueError: ``python-pptx`` 가 설치되어 있지 않을 때 발생.
        """
        try:
            from pptx import Presentation
        except ImportError as e:
            logger.warning("pptx_dependency_missing", error=str(e))
            raise ValueError(
                "python-pptx가 설치되어 있지 않습니다. 'uv sync'로 의존성을 설치하세요."
            ) from e

        try:
            presentation = Presentation(str(path))
            parts: list[str] = []

            for index, slide in enumerate(presentation.slides, start=1):
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_parts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            line = " | ".join(c for c in cells if c)
                            if line:
                                slide_parts.append(line)
                if slide_parts:
                    parts.append(f"# 슬라이드 {index}\n" + "\n".join(slide_parts))

            return "\n\n".join(parts)
        except Exception as e:
            logger.warning("pptx_parse_failed", path=str(path), error=str(e))
            return ""
